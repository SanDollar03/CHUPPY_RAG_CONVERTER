# app.py
# -*- coding: utf-8 -*-
import os
import re
import json
import time
import threading
import uuid
from collections import deque
from datetime import datetime
from typing import Dict, List, Tuple, Optional, Any

import requests
from flask import Flask, render_template, request, jsonify, Response
from dotenv import load_dotenv

from docx import Document
from pypdf import PdfReader

from openpyxl import load_workbook
import xlrd

from pptx import Presentation


load_dotenv()

APP_TITLE = "CHUPPY RAG CONVERTER"
HEADER_MODEL_LABEL = "Model : ChatGPT 5.2"

# Chat API (/chat-messages)
API_BASE = (os.getenv("DIFY_API_BASE") or "").strip().rstrip("/")
API_KEY = (os.getenv("DIFY_API_KEY") or "").strip()

# Knowledge API (/datasets ...)
DATASET_API_BASE = (os.getenv("DIFY_DATASET_API_BASE") or API_BASE).strip().rstrip("/")
DATASET_API_KEY = (os.getenv("DIFY_DATASET_API_KEY") or API_KEY).strip()
DATASET_NAME_PREFIX = (os.getenv("DATASET_NAME_PREFIX") or "Chu_").strip()

# On-demand explorer root
EXPLORER_ROOT = os.path.normpath(
    os.getenv("CHUPPY_EXPLORER_ROOT") or r"\\172.27.23.54\disk1\Chuppy"
)
UPLOAD_MAX_FILES = 100
EXPLORER_MAX_DEPTH = 5
UPLOAD_ALLOWED_DEPTH = 5

ALLOWED_EXTS = {
    ".txt", ".md", ".csv", ".json", ".log",
    ".html", ".xml", ".yml", ".yaml", ".ini", ".conf",
    ".py", ".js", ".css",
    ".docx", ".pdf",
    ".xlsx", ".xls", ".xlsm",
    ".ppt", ".pptx",
}

MAX_INPUT_CHARS = 180_000
DEFAULT_CHUNK_SEP = "***"
REQ_TIMEOUT_SEC = 300

INDEXING_POLL_SEC = float(os.getenv("DIFY_INDEXING_POLL_SEC") or "2.0")
INDEXING_MAX_WAIT_SEC = int(os.getenv("DIFY_INDEXING_MAX_WAIT_SEC") or "900")
DIFY_MAX_SEG_TOKENS = int(os.getenv("DIFY_MAX_SEG_TOKENS") or "2000")

ONDEMAND_QUEUE_MAX_RETRIES = int(os.getenv("ONDEMAND_QUEUE_MAX_RETRIES") or "5")
ONDEMAND_QUEUE_HISTORY_LIMIT = int(os.getenv("ONDEMAND_QUEUE_HISTORY_LIMIT") or "300")
ONDEMAND_DATASET_CACHE_TTL_SEC = int(os.getenv("ONDEMAND_DATASET_CACHE_TTL_SEC") or "60")
ONDEMAND_DOCUMENT_CACHE_TTL_SEC = int(os.getenv("ONDEMAND_DOCUMENT_CACHE_TTL_SEC") or "60")
ONDEMAND_QUEUE_USER = (os.getenv("ONDEMAND_QUEUE_USER") or "rag_converter").strip() or "rag_converter"
ONDEMAND_QUEUE_STYLE = "rag_markdown"
ONDEMAND_QUEUE_CHUNK_SEP = DEFAULT_CHUNK_SEP
ONDEMAND_MONITOR_ENABLED = str(os.getenv("ONDEMAND_MONITOR_ENABLED") or "1").strip().lower() not in {"0", "false", "no", "off"}
ONDEMAND_MONITOR_INTERVAL_SEC = max(3.0, float(os.getenv("ONDEMAND_MONITOR_INTERVAL_SEC") or "15"))
ONDEMAND_SEEN_SIGNATURE_LIMIT = max(1000, int(os.getenv("ONDEMAND_SEEN_SIGNATURE_LIMIT") or "5000"))

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
NOTICE_PATH = os.path.join(BASE_DIR, "notice.txt")


def ensure_notice_file() -> None:
    if os.path.exists(NOTICE_PATH):
        return
    try:
        with open(NOTICE_PATH, "w", encoding="utf-8") as f:
            f.write("")
    except Exception:
        pass


def create_app():
    ensure_notice_file()

    app = Flask(__name__)
    app.config["JSON_AS_ASCII"] = False

    ONDEMAND_QUEUE.start()
    ONDEMAND_MONITOR.start()

    @app.get("/")
    def index():
        return render_template(
            "index.html",
            title=APP_TITLE,
            model_label=HEADER_MODEL_LABEL,
            api_ready=bool(API_BASE and API_KEY),
        )

    @app.get("/auto")
    def auto_page():
        return render_template(
            "auto_rag.html",
            title=APP_TITLE + " (AUTO)",
            model_label=HEADER_MODEL_LABEL,
            api_ready=bool(API_BASE and API_KEY),
            dataset_api_ready=bool(DATASET_API_BASE and DATASET_API_KEY),
            dataset_prefix=DATASET_NAME_PREFIX,
        )

    @app.get("/knowledge")
    def knowledge_page():
        return render_template(
            "knowledge.html",
            title=APP_TITLE + " (KNOWLEDGE)",
            model_label=HEADER_MODEL_LABEL,
            dataset_api_ready=bool(DATASET_API_BASE and DATASET_API_KEY),
            dataset_prefix=DATASET_NAME_PREFIX,
        )

    @app.get("/ondemand")
    def ondemand_page():
        return render_template(
            "ondemand.html",
            title=APP_TITLE + " (ON DEMAND)",
            model_label=HEADER_MODEL_LABEL,
            explorer_root=EXPLORER_ROOT,
            explorer_max_depth=EXPLORER_MAX_DEPTH,
            upload_allowed_depth=UPLOAD_ALLOWED_DEPTH,
        )

    @app.get("/api/health")
    def api_health():
        return jsonify({
            "ok": True,
            "api_ready": bool(API_BASE and API_KEY),
            "dataset_api_ready": bool(DATASET_API_BASE and DATASET_API_KEY),
            "dataset_prefix": DATASET_NAME_PREFIX,
            "model_label": HEADER_MODEL_LABEL,
            "explorer_root": EXPLORER_ROOT,
            "explorer_max_depth": EXPLORER_MAX_DEPTH,
            "upload_allowed_depth": UPLOAD_ALLOWED_DEPTH,
        })

    @app.get("/api/notice")
    def api_notice():
        ensure_notice_file()
        try:
            with open(NOTICE_PATH, "r", encoding="utf-8", errors="ignore") as f:
                txt = f.read()
        except Exception:
            return jsonify({"ok": False, "error": "notice.txt の読み取りに失敗しました。"}), 500

        if len(txt) > 50_000:
            txt = txt[:50_000] + "\n...(truncated)\n"

        return jsonify({"ok": True, "text": txt})

    @app.get("/api/datasets")
    def api_datasets():
        if not DATASET_API_BASE or not DATASET_API_KEY:
            return jsonify({"ok": False, "error": "ナレッジAPI設定が未完了です（DIFY_DATASET_API_BASE / DIFY_DATASET_API_KEY）。"}), 500

        try:
            items = dify_list_datasets(
                api_base=DATASET_API_BASE,
                api_key=DATASET_API_KEY,
                prefix=DATASET_NAME_PREFIX,
                limit=100,
            )
        except Exception as e:
            return jsonify({"ok": False, "error": safe_err(str(e))}), 500

        return jsonify({"ok": True, "items": items, "prefix": DATASET_NAME_PREFIX})

    @app.get("/api/knowledge/datasets/<dataset_id>/detail")
    def api_knowledge_dataset_detail(dataset_id):
        if not DATASET_API_BASE or not DATASET_API_KEY:
            return jsonify({"ok": False, "error": "ナレッジAPI設定が未完了です（DIFY_DATASET_API_BASE / DIFY_DATASET_API_KEY）。"}), 500

        try:
            ds = dify_get_dataset_detail(
                api_base=DATASET_API_BASE,
                api_key=DATASET_API_KEY,
                dataset_id=dataset_id,
            )
        except Exception as e:
            return jsonify({"ok": False, "error": safe_err(str(e))}), 500

        out = {"ok": True}
        if isinstance(ds, dict):
            out.update(ds)
        else:
            out["item"] = ds
        return jsonify(out)

    @app.get("/api/knowledge/datasets/<dataset_id>/documents")
    def api_knowledge_documents(dataset_id):
        if not DATASET_API_BASE or not DATASET_API_KEY:
            return jsonify({"ok": False, "error": "ナレッジAPI設定が未完了です（DIFY_DATASET_API_BASE / DIFY_DATASET_API_KEY）。"}), 500

        keyword = (request.args.get("keyword") or "").strip()

        try:
            items, total = dify_list_documents_all(
                api_base=DATASET_API_BASE,
                api_key=DATASET_API_KEY,
                dataset_id=dataset_id,
                keyword=keyword,
                limit=100,
            )
        except Exception as e:
            return jsonify({"ok": False, "error": safe_err(str(e))}), 500

        return jsonify({"ok": True, "items": items, "total": total, "keyword": keyword})

    @app.get("/api/knowledge/datasets/<dataset_id>/documents/<document_id>")
    def api_knowledge_document_detail(dataset_id, document_id):
        if not DATASET_API_BASE or not DATASET_API_KEY:
            return jsonify({"ok": False, "error": "ナレッジAPI設定が未完了です（DIFY_DATASET_API_BASE / DIFY_DATASET_API_KEY）。"}), 500

        metadata = (request.args.get("metadata") or "without").strip() or "without"

        try:
            doc = dify_get_document_detail(
                api_base=DATASET_API_BASE,
                api_key=DATASET_API_KEY,
                dataset_id=dataset_id,
                document_id=document_id,
                metadata=metadata,
            )
        except Exception as e:
            return jsonify({"ok": False, "error": safe_err(str(e))}), 500

        out = {"ok": True}
        if isinstance(doc, dict):
            out.update(doc)
        else:
            out["item"] = doc
        return jsonify(out)

    @app.get("/api/knowledge/datasets/<dataset_id>/documents/<document_id>/segments")
    def api_knowledge_segments(dataset_id, document_id):
        if not DATASET_API_BASE or not DATASET_API_KEY:
            return jsonify({"ok": False, "error": "ナレッジAPI設定が未完了です（DIFY_DATASET_API_BASE / DIFY_DATASET_API_KEY）。"}), 500

        keyword = (request.args.get("keyword") or "").strip()
        status = (request.args.get("status") or "").strip()

        try:
            page = int(request.args.get("page") or "1")
            limit = int(request.args.get("limit") or "20")
        except Exception:
            page = 1
            limit = 20

        page = max(1, page)
        limit = max(1, min(100, limit))

        try:
            res = dify_list_segments_page(
                api_base=DATASET_API_BASE,
                api_key=DATASET_API_KEY,
                dataset_id=dataset_id,
                document_id=document_id,
                page=page,
                limit=limit,
                keyword=keyword,
                status=status,
            )
        except Exception as e:
            return jsonify({"ok": False, "error": safe_err(str(e))}), 500

        return jsonify({
            "ok": True,
            "items": res.get("items") or [],
            "has_more": bool(res.get("has_more")),
            "total": int(res.get("total") or 0),
            "page": int(res.get("page") or page),
            "limit": int(res.get("limit") or limit),
        })

    @app.get("/api/knowledge/datasets/<dataset_id>/documents/<document_id>/segments/<segment_id>")
    def api_knowledge_segment_detail(dataset_id, document_id, segment_id):
        if not DATASET_API_BASE or not DATASET_API_KEY:
            return jsonify({"ok": False, "error": "ナレッジAPI設定が未完了です（DIFY_DATASET_API_BASE / DIFY_DATASET_API_KEY）。"}), 500

        try:
            seg = dify_get_segment_detail(
                api_base=DATASET_API_BASE,
                api_key=DATASET_API_KEY,
                dataset_id=dataset_id,
                document_id=document_id,
                segment_id=segment_id,
            )
        except Exception as e:
            return jsonify({"ok": False, "error": safe_err(str(e))}), 500

        out = {"ok": True}
        if isinstance(seg, dict):
            out.update(seg)
        else:
            out["item"] = seg
        return jsonify(out)

    @app.get("/api/explorer/root")
    def api_explorer_root():
        try:
            stats_cache: Dict[str, Dict[str, int]] = {}
            root_info = build_dir_info(EXPLORER_ROOT, EXPLORER_ROOT, stats_cache)
            root_info["depth"] = 0
            root_info["can_upload"] = False
            root_info["has_children"] = dir_has_child_dirs(EXPLORER_ROOT)
            return jsonify({
                "ok": True,
                "root": root_info,
                "max_depth": EXPLORER_MAX_DEPTH,
                "upload_allowed_depth": UPLOAD_ALLOWED_DEPTH,
            })
        except Exception as e:
            return jsonify({"ok": False, "error": safe_err(str(e))}), 500

    @app.get("/api/explorer/list")
    def api_explorer_list():
        rel_path = (request.args.get("path") or "").strip()

        try:
            abs_dir = resolve_explorer_path(EXPLORER_ROOT, rel_path)
            if not os.path.isdir(abs_dir):
                return jsonify({"ok": False, "error": "対象フォルダが存在しません。"}), 404

            dirs, files = list_explorer_dir(abs_dir, EXPLORER_ROOT)
            stats_cache: Dict[str, Dict[str, int]] = {}
            current = build_dir_info(abs_dir, EXPLORER_ROOT, stats_cache)
            return jsonify({
                "ok": True,
                "current": current,
                "dirs": dirs,
                "files": files,
                "max_depth": EXPLORER_MAX_DEPTH,
                "upload_allowed_depth": UPLOAD_ALLOWED_DEPTH,
            })
        except Exception as e:
            return jsonify({"ok": False, "error": safe_err(str(e))}), 400

    @app.post("/api/explorer/upload")
    def api_explorer_upload():
        rel_path = (request.form.get("path") or "").strip()

        try:
            target_dir = resolve_explorer_path(EXPLORER_ROOT, rel_path)
        except Exception as e:
            return jsonify({"ok": False, "error": safe_err(str(e))}), 400

        if not os.path.isdir(target_dir):
            return jsonify({"ok": False, "error": "保存先フォルダが存在しません。"}), 404

        depth = path_depth_from_root(target_dir, EXPLORER_ROOT)
        if depth != UPLOAD_ALLOWED_DEPTH:
            return jsonify({
                "ok": False,
                "error": f"ファイル追加が許可されているのは {UPLOAD_ALLOWED_DEPTH} 階層目のフォルダのみです。"
            }), 400

        files = request.files.getlist("files")
        if not files:
            return jsonify({"ok": False, "error": "アップロードするファイルがありません。"}), 400

        if len(files) > UPLOAD_MAX_FILES:
            return jsonify({"ok": False, "error": f"一度にアップロードできるのは最大 {UPLOAD_MAX_FILES} 件です。"}), 400

        saved = []
        skipped = []
        errors = []
        queue_items = []
        queue_errors = []
        folder_rel_path = make_rel_from_root(target_dir, EXPLORER_ROOT)

        for f in files:
            try:
                original_name_raw = (f.filename or "").strip()
                if not original_name_raw:
                    skipped.append({"name": "", "reason": "ファイル名が空です。"})
                    continue

                original_name = sanitize_upload_filename(original_name_raw)
                if not original_name:
                    skipped.append({"name": original_name_raw, "reason": "使用できないファイル名です。"})
                    continue

                stored_name = add_upload_timestamp_prefix(original_name)
                save_path = build_unique_upload_path(target_dir, stored_name)
                saved_name = os.path.basename(save_path)
                f.save(save_path)

                saved_item = {
                    "name": saved_name,
                    "original_name": original_name,
                    "rel_path": make_rel_from_root(save_path, EXPLORER_ROOT),
                    "size_bytes": os.path.getsize(save_path) if os.path.exists(save_path) else 0,
                }
                saved.append(saved_item)

                try:
                    task = ONDEMAND_QUEUE.enqueue_saved_file(
                        folder_rel_path=folder_rel_path,
                        folder_abs_path=target_dir,
                        source_abs_path=save_path,
                        source_rel_path=saved_item["rel_path"],
                        source_saved_name=saved_name,
                        source_original_name=original_name,
                    )
                    if task:
                        queue_items.append(task)
                except Exception as qerr:
                    queue_errors.append({
                        "name": original_name,
                        "error": safe_err(str(qerr)),
                    })
            except Exception as e:
                errors.append({"name": f.filename or "", "error": safe_err(str(e))})

        return jsonify({
            "ok": True,
            "target": build_dir_info(target_dir, EXPLORER_ROOT, {}),
            "saved": saved,
            "skipped": skipped,
            "errors": errors,
            "queue_items": queue_items,
            "queue_errors": queue_errors,
        })

    @app.get("/api/ondemand/queue")
    def api_ondemand_queue():
        try:
            limit = int(request.args.get("limit") or "200")
        except Exception:
            limit = 200
        limit = max(20, min(500, limit))

        try:
            snap = ONDEMAND_QUEUE.get_snapshot(limit=limit)
            monitor = ONDEMAND_MONITOR.get_status()
            return jsonify({"ok": True, **snap, "monitor": monitor})
        except Exception as e:
            return jsonify({"ok": False, "error": safe_err(str(e))}), 500

    @app.post("/api/scan")
    def api_scan():
        data = request.get_json(force=True) or {}
        in_dir = (data.get("input_dir") or "").strip()
        recursive = bool(data.get("recursive", True))

        if not in_dir or not os.path.isdir(in_dir):
            return jsonify({"ok": False, "error": "入力フォルダが存在しません。"}), 400

        files = list_files(in_dir, recursive=recursive)
        return jsonify({"ok": True, "count": len(files), "files": files})

    @app.post("/api/run")
    def api_run():
        if not API_BASE or not API_KEY:
            return jsonify({
                "ok": False,
                "error": "サーバー側API設定が未完了です。.env に DIFY_API_BASE / DIFY_API_KEY を設定してください。"
            }), 500

        data = request.get_json(force=True) or {}

        input_dir = (data.get("input_dir") or "").strip()
        output_dir = (data.get("output_dir") or "").strip()
        recursive = bool(data.get("recursive", True))

        user = (data.get("user") or "rag_converter").strip()
        knowledge_style = (data.get("knowledge_style") or "rag_markdown").strip()
        chunk_sep = (data.get("chunk_sep") or DEFAULT_CHUNK_SEP).strip() or DEFAULT_CHUNK_SEP

        overwrite = bool(data.get("overwrite", False))

        if not input_dir or not os.path.isdir(input_dir):
            return jsonify({"ok": False, "error": "入力フォルダが存在しません。"}), 400
        if not output_dir:
            return jsonify({"ok": False, "error": "出力フォルダが未指定です。"}), 400

        os.makedirs(output_dir, exist_ok=True)
        files = list_files(input_dir, recursive=recursive)

        def sse():
            yield sse_event("meta", {
                "title": APP_TITLE,
                "mode": "manual",
                "model": HEADER_MODEL_LABEL,
                "total": len(files),
                "overwrite": overwrite,
            })

            ok_count = 0
            ng_count = 0
            skip_count = 0

            for idx, relpath in enumerate(files, start=1):
                abspath = os.path.join(input_dir, relpath)
                yield sse_event("progress", {"index": idx, "total": len(files), "file": relpath})

                try:
                    out_path = make_output_path(output_dir, relpath)

                    if (not overwrite) and os.path.exists(out_path):
                        skip_count += 1
                        yield sse_event("skip_one", {"file": relpath, "out": os.path.relpath(out_path, output_dir)})
                        continue

                    raw_text, meta = extract_text(abspath, knowledge_style=knowledge_style)

                    if not raw_text.strip():
                        raise RuntimeError("抽出テキストが空でした。")

                    if len(raw_text) > MAX_INPUT_CHARS:
                        raw_text = raw_text[:MAX_INPUT_CHARS] + "\n...(truncated)\n"

                    md_body = convert_via_dify_chat_messages_secure(
                        api_base=API_BASE,
                        api_key=API_KEY,
                        user=user,
                        source_path=relpath,
                        source_meta=meta,
                        text=raw_text,
                        knowledge_style=knowledge_style,
                        chunk_sep=chunk_sep,
                    )

                    md_body = normalize_chunk_sep_lines(md_body, chunk_sep)

                    md_save = attach_source_metadata(
                        md_body,
                        source_relpath=relpath,
                        source_abspath=abspath,
                        source_meta=meta,
                    )

                    os.makedirs(os.path.dirname(out_path), exist_ok=True)
                    with open(out_path, "w", encoding="utf-8", newline="\n") as f:
                        f.write(md_save)

                    ok_count += 1
                    yield sse_event("done_one", {"file": relpath, "out": os.path.relpath(out_path, output_dir)})

                except Exception as e:
                    ng_count += 1
                    yield sse_event("error_one", {"file": relpath, "error": safe_err(str(e))})

            yield sse_event("summary", {
                "ok": ok_count,
                "ng": ng_count,
                "skip": skip_count,
                "total": len(files),
                "overwrite": overwrite,
            })

        return Response(sse(), mimetype="text/event-stream")

    @app.post("/api/auto/run")
    def api_auto_run():
        if not API_BASE or not API_KEY:
            return jsonify({
                "ok": False,
                "error": "生成AI API設定が未完了です。.env に DIFY_API_BASE / DIFY_API_KEY を設定してください。"
            }), 500

        if not DATASET_API_BASE or not DATASET_API_KEY:
            return jsonify({
                "ok": False,
                "error": "ナレッジAPI設定が未完了です。.env に DIFY_DATASET_API_BASE / DIFY_DATASET_API_KEY を設定してください。"
            }), 500

        data = request.get_json(force=True) or {}

        input_dir = (data.get("input_dir") or "").strip()
        output_dir = (data.get("output_dir") or "").strip()
        recursive = bool(data.get("recursive", True))

        dataset_id = (data.get("dataset_id") or "").strip()

        user = (data.get("user") or "rag_converter").strip()
        knowledge_style = (data.get("knowledge_style") or "rag_markdown").strip()
        chunk_sep = (data.get("chunk_sep") or DEFAULT_CHUNK_SEP).strip() or DEFAULT_CHUNK_SEP

        overwrite = bool(data.get("overwrite", False))

        if not dataset_id:
            return jsonify({"ok": False, "error": "ナレッジ（dataset_id）が未指定です。"}), 400

        if not input_dir or not os.path.isdir(input_dir):
            return jsonify({"ok": False, "error": "入力フォルダが存在しません。"}), 400
        if not output_dir:
            return jsonify({"ok": False, "error": "出力フォルダが未指定です。"}), 400

        os.makedirs(output_dir, exist_ok=True)
        files = list_files(input_dir, recursive=recursive)

        def sse():
            yield sse_event("meta", {
                "title": APP_TITLE,
                "mode": "auto",
                "model": HEADER_MODEL_LABEL,
                "total": len(files),
                "overwrite": overwrite,
                "dataset_id": dataset_id,
                "chunk_sep": chunk_sep,
            })

            ok_count = 0
            ng_count = 0
            skip_count = 0

            for idx, relpath in enumerate(files, start=1):
                abspath = os.path.join(input_dir, relpath)
                yield sse_event("progress", {"index": idx, "total": len(files), "file": relpath})

                try:
                    out_path = make_output_path(output_dir, relpath)

                    if (not overwrite) and os.path.exists(out_path):
                        skip_count += 1
                        yield sse_event("skip_one", {"file": relpath, "out": os.path.relpath(out_path, output_dir)})
                        continue

                    raw_text, meta = extract_text(abspath, knowledge_style=knowledge_style)

                    if not raw_text.strip():
                        raise RuntimeError("抽出テキストが空でした。")

                    if len(raw_text) > MAX_INPUT_CHARS:
                        raw_text = raw_text[:MAX_INPUT_CHARS] + "\n...(truncated)\n"

                    md_body = convert_via_dify_chat_messages_secure(
                        api_base=API_BASE,
                        api_key=API_KEY,
                        user=user,
                        source_path=relpath,
                        source_meta=meta,
                        text=raw_text,
                        knowledge_style=knowledge_style,
                        chunk_sep=chunk_sep,
                    )

                    md_body = normalize_chunk_sep_lines(md_body, chunk_sep)

                    md_save = attach_source_metadata(
                        md_body,
                        source_relpath=relpath,
                        source_abspath=abspath,
                        source_meta=meta,
                    )

                    os.makedirs(os.path.dirname(out_path), exist_ok=True)
                    with open(out_path, "w", encoding="utf-8", newline="\n") as f:
                        f.write(md_save)

                    yield sse_event("done_one", {"file": relpath, "out": os.path.relpath(out_path, output_dir)})

                    reg = register_markdown_to_dify(
                        dataset_id=dataset_id,
                        doc_name=os.path.basename(out_path),
                        markdown=md_body,
                        chunk_sep=chunk_sep,
                    )

                    yield sse_event("dataset", {
                        "file": relpath,
                        "doc_id": reg.get("doc_id"),
                        "batch": reg.get("batch"),
                        "chunk_sep": reg.get("chunk_sep"),
                        "chunks": reg.get("chunks"),
                        "chunk_tokens_max": reg.get("chunk_tokens_max"),
                        "dify_max_tokens": reg.get("dify_max_tokens"),
                        "search_method": reg.get("search_method"),
                        "message": "ナレッジ登録 受付",
                    })

                    final = None
                    for prog in iter_indexing_status(
                        dataset_id=dataset_id,
                        batch=reg["batch"],
                        doc_id=reg["doc_id"],
                    ):
                        yield sse_event("dataset_progress", {
                            "file": relpath,
                            "doc_id": reg.get("doc_id"),
                            "batch": reg.get("batch"),
                            "status": prog.get("indexing_status"),
                            "completed_segments": prog.get("completed_segments"),
                            "total_segments": prog.get("total_segments"),
                            "error": prog.get("error"),
                            "terminal": bool(prog.get("terminal")),
                        })
                        if prog.get("terminal"):
                            final = prog
                            break

                    if not final:
                        raise RuntimeError("Dify埋め込みの進捗取得に失敗しました。")

                    if (final.get("indexing_status") or "").lower() != "completed":
                        raise RuntimeError(
                            f"Dify埋め込み失敗: status={final.get('indexing_status')} error={final.get('error')}"
                        )

                    if int(final.get("total_segments") or 0) <= 0:
                        raise RuntimeError("Dify側で0セグメントのまま完了しました（separator/max_tokens/text を要確認）。")

                    yield sse_event("dataset_done", {
                        "file": relpath,
                        "doc_id": reg.get("doc_id"),
                        "batch": reg.get("batch"),
                        "status": final.get("indexing_status"),
                        "completed_segments": final.get("completed_segments"),
                        "total_segments": final.get("total_segments"),
                        "message": "ナレッジ登録 完了",
                    })

                    ok_count += 1

                except Exception as e:
                    ng_count += 1
                    yield sse_event("error_one", {"file": relpath, "error": safe_err(str(e))})

            yield sse_event("summary", {
                "ok": ok_count,
                "ng": ng_count,
                "skip": skip_count,
                "total": len(files),
                "overwrite": overwrite,
            })

        return Response(sse(), mimetype="text/event-stream")

    return app


# -----------------------
# File utilities
# -----------------------

def list_files(root_dir: str, recursive: bool = True) -> List[str]:
    results: List[str] = []
    root_dir = os.path.abspath(root_dir)

    if recursive:
        for base, _, files in os.walk(root_dir):
            for name in files:
                ext = os.path.splitext(name)[1].lower()
                if ext in ALLOWED_EXTS:
                    abs_path = os.path.join(base, name)
                    rel = os.path.relpath(abs_path, root_dir)
                    results.append(rel)
    else:
        for name in os.listdir(root_dir):
            abs_path = os.path.join(root_dir, name)
            if os.path.isfile(abs_path):
                ext = os.path.splitext(name)[1].lower()
                if ext in ALLOWED_EXTS:
                    results.append(name)

    results.sort()
    return results


def extract_text(path: str, knowledge_style: str = "rag_markdown") -> Tuple[str, Dict[str, str]]:
    ext = os.path.splitext(path)[1].lower()
    stat = os.stat(path)
    meta = {
        "filename": os.path.basename(path),
        "ext": ext,
        "size_bytes": str(stat.st_size),
        "mtime": datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d %H:%M:%S"),
    }

    if ext in {
        ".txt", ".md", ".csv", ".json", ".log",
        ".html", ".xml", ".yml", ".yaml", ".ini", ".conf",
        ".py", ".js", ".css",
    }:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read(), meta

    if ext == ".docx":
        doc = Document(path)
        parts = []
        for p in doc.paragraphs:
            t = p.text.strip()
            if t:
                parts.append(t)
        return "\n".join(parts), meta

    if ext == ".pdf":
        return extract_pdf_like(path), meta

    if ext in {".xlsx", ".xlsm", ".xls"}:
        if knowledge_style == "rag_natural":
            text = extract_excel_as_markdown_tables(path, ext)
        else:
            text = extract_excel_as_row_records(path, ext)
        return text, meta

    if ext in {".ppt", ".pptx"}:
        return extract_ppt_like(path, ext), meta

    raise RuntimeError(f"未対応の拡張子です: {ext}")


def extract_pdf_like(path: str) -> str:
    reader = PdfReader(path)
    parts = []
    for i, page in enumerate(reader.pages):
        txt = page.extract_text() or ""
        txt = normalize_pdf_like_text(txt)
        if txt.strip():
            parts.append(f"[PAGE {i+1}]\n{txt}")
    return "\n\n".join(parts)


def extract_excel_as_row_records(path: str, ext: str) -> str:
    if ext == ".xls":
        return extract_xls_as_row_records(path)
    return extract_xlsx_like_as_row_records(path)


def extract_xlsx_like_as_row_records(path: str) -> str:
    wb = load_workbook(path, data_only=True, read_only=True)
    out: List[str] = []

    for sheet in wb.worksheets:
        out.append(f"[SHEET: {sheet.title}]")

        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            out.append("[EMPTY]")
            out.append("")
            continue

        header: Optional[List[str]] = None
        start_idx = 0
        for i, r in enumerate(rows):
            if any(cell is not None and str(cell).strip() != "" for cell in r):
                header = [sanitize_header(cell) for cell in r]
                start_idx = i + 1
                break

        if not header:
            out.append("[EMPTY]")
            out.append("")
            continue

        out.append("[HEADER] " + "\t".join([h if h else "" for h in header]))

        for ridx in range(start_idx, len(rows)):
            r = rows[ridx]
            if not any(cell is not None and str(cell).strip() != "" for cell in r):
                continue

            record: Dict[str, str] = {}
            for cidx, cell in enumerate(r):
                key = header[cidx] if cidx < len(header) else f"COL{cidx+1}"
                if not key:
                    key = f"COL{cidx+1}"
                val = "" if cell is None else str(cell).strip()
                if val != "":
                    record[key] = val

            if record:
                out.append("[ROW] " + json.dumps(record, ensure_ascii=False, separators=(",", ":")))

        out.append("")

    return "\n".join(out).strip()


def extract_xls_as_row_records(path: str) -> str:
    wb = xlrd.open_workbook(path)
    out: List[str] = []

    for sheet in wb.sheets():
        out.append(f"[SHEET: {sheet.name}]")

        if sheet.nrows <= 0:
            out.append("[EMPTY]")
            out.append("")
            continue

        rows = []
        for r in range(sheet.nrows):
            rows.append([sheet.cell_value(r, c) for c in range(sheet.ncols)])

        header: Optional[List[str]] = None
        start_idx = 0
        for i, r in enumerate(rows):
            if any(cell is not None and str(cell).strip() != "" for cell in r):
                header = [sanitize_header(cell) for cell in r]
                start_idx = i + 1
                break

        if not header:
            out.append("[EMPTY]")
            out.append("")
            continue

        out.append("[HEADER] " + "\t".join([h if h else "" for h in header]))

        for ridx in range(start_idx, len(rows)):
            r = rows[ridx]
            if not any(cell is not None and str(cell).strip() != "" for cell in r):
                continue

            record: Dict[str, str] = {}
            for cidx, cell in enumerate(r):
                key = header[cidx] if cidx < len(header) else f"COL{cidx+1}"
                if not key:
                    key = f"COL{cidx+1}"
                val = "" if cell is None else str(cell).strip()
                if val != "":
                    record[key] = val

            if record:
                out.append("[ROW] " + json.dumps(record, ensure_ascii=False, separators=(",", ":")))

        out.append("")

    return "\n".join(out).strip()


def extract_excel_as_markdown_tables(path: str, ext: str) -> str:
    if ext == ".xls":
        return extract_xls_as_markdown_tables(path)
    return extract_xlsx_like_as_markdown_tables(path)


def extract_xlsx_like_as_markdown_tables(path: str) -> str:
    max_rows_per_sheet = 200
    wb = load_workbook(path, data_only=True, read_only=True)

    out: List[str] = []
    for sheet in wb.worksheets:
        out.append(f"[SHEET: {sheet.title}]")

        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            out.append("(empty)")
            out.append("")
            continue

        header = None
        start_idx = 0
        for i, r in enumerate(rows):
            if any(cell is not None and str(cell).strip() != "" for cell in r):
                header = [sanitize_header(c) for c in r]
                start_idx = i + 1
                break
        if not header:
            out.append("(empty)")
            out.append("")
            continue

        cols = [h if h else f"COL{j+1}" for j, h in enumerate(header)]

        out.append("| " + " | ".join(cols) + " |")
        out.append("| " + " | ".join(["---"] * len(cols)) + " |")

        count = 0
        for ridx in range(start_idx, len(rows)):
            r = rows[ridx]
            if not any(cell is not None and str(cell).strip() != "" for cell in r):
                continue

            vals = []
            for cidx in range(len(cols)):
                cell = r[cidx] if cidx < len(r) else None
                v = "" if cell is None else str(cell).strip()
                v = v.replace("\n", " ").replace("\r", " ")
                v = v.replace("|", "\\|")
                vals.append(v)

            out.append("| " + " | ".join(vals) + " |")
            count += 1
            if count >= max_rows_per_sheet:
                out.append(f"(… {max_rows_per_sheet}行まで表示。続きは省略 …)")
                break

        out.append("")

    return "\n".join(out).strip()


def extract_xls_as_markdown_tables(path: str) -> str:
    max_rows_per_sheet = 200
    wb = xlrd.open_workbook(path)
    out: List[str] = []

    for sheet in wb.sheets():
        out.append(f"[SHEET: {sheet.name}]")

        if sheet.nrows <= 0:
            out.append("(empty)")
            out.append("")
            continue

        rows = []
        for r in range(sheet.nrows):
            rows.append([sheet.cell_value(r, c) for c in range(sheet.ncols)])

        header = None
        start_idx = 0
        for i, r in enumerate(rows):
            if any(cell is not None and str(cell).strip() != "" for cell in r):
                header = [sanitize_header(c) for c in r]
                start_idx = i + 1
                break
        if not header:
            out.append("(empty)")
            out.append("")
            continue

        cols = [h if h else f"COL{j+1}" for j, h in enumerate(header)]

        out.append("| " + " | ".join(cols) + " |")
        out.append("| " + " | ".join(["---"] * len(cols)) + " |")

        count = 0
        for ridx in range(start_idx, len(rows)):
            r = rows[ridx]
            if not any(cell is not None and str(cell).strip() != "" for cell in r):
                continue

            vals = []
            for cidx in range(len(cols)):
                cell = r[cidx] if cidx < len(r) else None
                v = "" if cell is None else str(cell).strip()
                v = v.replace("\n", " ").replace("\r", " ")
                v = v.replace("|", "\\|")
                vals.append(v)

            out.append("| " + " | ".join(vals) + " |")
            count += 1
            if count >= max_rows_per_sheet:
                out.append(f"(… {max_rows_per_sheet}行まで表示。続きは省略 …)")
                break

        out.append("")

    return "\n".join(out).strip()


def extract_ppt_like(path: str, ext: str) -> str:
    try:
        prs = Presentation(path)
    except Exception:
        if ext == ".ppt":
            raise RuntimeError("`.ppt`（旧形式）は python-pptx で直接読めない場合があります。`.pptx` に変換して再実行してください。")
        raise RuntimeError("PowerPointの解析に失敗しました。ファイル破損または形式が想定外です。")

    parts: List[str] = []
    for i, slide in enumerate(prs.slides):
        slide_text: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = (shape.text or "").strip()
                if t:
                    slide_text.append(t)

        txt = "\n".join(slide_text)
        txt = normalize_pdf_like_text(txt)
        if txt.strip():
            parts.append(f"[SLIDE {i+1}]\n{txt}")

    return "\n\n".join(parts)


def sanitize_header(cell) -> str:
    if cell is None:
        return ""
    s = str(cell).strip()
    s = re.sub(r"\s+", " ", s)
    return s


def normalize_pdf_like_text(s: str) -> str:
    lines = [ln.rstrip() for ln in s.splitlines()]
    out: List[str] = []
    buf = ""

    def flush():
        nonlocal buf
        if buf:
            out.append(buf)
            buf = ""

    for ln in lines:
        t = ln.strip("\u00a0 ").strip()
        if not t:
            flush()
            out.append("")
            continue
        if len(t) == 1:
            buf += t
        else:
            flush()
            out.append(t)
    flush()

    text = "\n".join(out)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def make_output_path(output_dir: str, rel_input_path: str) -> str:
    rel_dir = os.path.dirname(rel_input_path)
    base_name = os.path.splitext(os.path.basename(rel_input_path))[0]
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    out_name = f"{ts}_{base_name}.md"

    safe_dir = sanitize_relpath(rel_dir) if rel_dir else ""
    return os.path.join(output_dir, safe_dir, out_name)


def sanitize_relpath(p: str) -> str:
    if not p:
        return ""
    p = p.replace("..", "__")
    p = re.sub(r'[<>:"|?*]', "_", p)
    return p


def normalize_root_path(p: str) -> str:
    if not p:
        raise RuntimeError("ルートパスが未設定です。")
    return os.path.normcase(os.path.abspath(os.path.normpath(p)))


def make_rel_from_root(abs_path: str, root_dir: str) -> str:
    rel = os.path.relpath(abs_path, root_dir)
    if rel == ".":
        return ""
    return rel.replace("\\", "/")


def path_depth_from_rel(rel_path: str) -> int:
    rel = (rel_path or "").strip().replace("\\", "/").strip("/")
    if not rel:
        return 0
    return len([p for p in rel.split("/") if p])


def path_depth_from_root(abs_path: str, root_dir: str) -> int:
    rel = make_rel_from_root(abs_path, root_dir)
    return path_depth_from_rel(rel)


def resolve_explorer_path(root_dir: str, rel_path: str) -> str:
    root_norm = normalize_root_path(root_dir)

    rel = (rel_path or "").strip().replace("/", os.sep).replace("\\", os.sep)
    rel = rel.lstrip(os.sep)

    candidate = os.path.normpath(os.path.join(root_dir, rel))
    cand_norm = normalize_root_path(candidate)

    if cand_norm != root_norm and not cand_norm.startswith(root_norm + os.sep):
        raise RuntimeError("許可されていないパスです。")

    return candidate


def sanitize_upload_filename(name: str) -> str:
    original = os.path.basename((name or "").replace("\x00", "").strip())
    if not original or original in {".", ".."}:
        return ""

    safe = re.sub(r"[\x00-\x1f]", "", original)
    safe = safe.replace("/", "_").replace("\\", "_")
    safe = re.sub(r'[:*?"<>|]', "_", safe)
    safe = safe.rstrip(" .")

    if not safe or safe in {".", ".."}:
        return ""
    return safe


def add_upload_timestamp_prefix(filename: str) -> str:
    safe = sanitize_upload_filename(filename)
    if not safe:
        return ""
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    return f"{ts}_{safe}"


def build_unique_upload_path(target_dir: str, filename: str) -> str:
    safe = sanitize_upload_filename(filename)
    if not safe:
        raise RuntimeError("使用できないファイル名です。")

    base, ext = os.path.splitext(safe)
    candidate = os.path.join(target_dir, safe)
    seq = 1
    while os.path.exists(candidate):
        candidate = os.path.join(target_dir, f"{base}_{seq:02d}{ext}")
        seq += 1
        if seq > 9999:
            raise RuntimeError("同名ファイルが多すぎるため保存できません。")
    return candidate


def matches_explorer_level_rule(depth: int, name: str) -> bool:
    nm = str(name or "")
    if depth == 1:
        return len(nm) == 1
    if depth == 2:
        return len(nm) == 2
    if depth == 4:
        return nm == "元データ"
    return True


def list_visible_child_dir_names(abs_dir: str, root_dir: str) -> List[str]:
    parent_depth = path_depth_from_root(abs_dir, root_dir)
    next_depth = parent_depth + 1
    if next_depth > EXPLORER_MAX_DEPTH:
        return []

    names: List[str] = []
    try:
        for name in os.listdir(abs_dir):
            full = os.path.join(abs_dir, name)
            if not os.path.isdir(full):
                continue
            if not matches_explorer_level_rule(next_depth, name):
                continue
            names.append(name)
    except Exception:
        return []

    names.sort(key=lambda x: x.lower())
    return names


def dir_has_child_dirs(abs_dir: str) -> bool:
    try:
        for name in os.listdir(abs_dir):
            full = os.path.join(abs_dir, name)
            if os.path.isdir(full):
                return True
    except Exception:
        return False
    return False


def dir_has_visible_child_dirs(abs_dir: str, root_dir: str) -> bool:
    return bool(list_visible_child_dir_names(abs_dir, root_dir))


def compute_visible_tree_stats(abs_dir: str, root_dir: str, cache: Optional[Dict[str, Dict[str, int]]] = None) -> Dict[str, int]:
    cache = cache if cache is not None else {}
    key = normalize_root_path(abs_dir)
    if key in cache:
        return cache[key]

    depth = path_depth_from_root(abs_dir, root_dir)
    child_names = list_visible_child_dir_names(abs_dir, root_dir)

    file_count = 0
    total_size_bytes = 0

    if depth >= EXPLORER_MAX_DEPTH or not child_names:
        try:
            for name in os.listdir(abs_dir):
                full = os.path.join(abs_dir, name)
                if not os.path.isfile(full):
                    continue
                try:
                    st = os.stat(full)
                except Exception:
                    continue
                file_count += 1
                total_size_bytes += int(st.st_size or 0)
        except Exception:
            pass
    else:
        for name in child_names:
            full = os.path.join(abs_dir, name)
            child_stats = compute_visible_tree_stats(full, root_dir, cache)
            file_count += int(child_stats.get("file_count") or 0)
            total_size_bytes += int(child_stats.get("total_size_bytes") or 0)

    out = {
        "file_count": file_count,
        "total_size_bytes": total_size_bytes,
    }
    cache[key] = out
    return out


def build_dir_info(abs_dir: str, root_dir: str, stats_cache: Optional[Dict[str, Dict[str, int]]] = None) -> Dict[str, Any]:
    depth = path_depth_from_root(abs_dir, root_dir)
    stats = compute_visible_tree_stats(abs_dir, root_dir, stats_cache)
    return {
        "name": os.path.basename(abs_dir.rstrip("\\/")) or abs_dir,
        "path": make_rel_from_root(abs_dir, root_dir),
        "abs_path": abs_dir,
        "depth": depth,
        "can_upload": depth == UPLOAD_ALLOWED_DEPTH,
        "has_children": dir_has_visible_child_dirs(abs_dir, root_dir),
        "file_count": int(stats.get("file_count") or 0),
        "total_size_bytes": int(stats.get("total_size_bytes") or 0),
    }


def list_explorer_dir(abs_dir: str, root_dir: str) -> Tuple[List[Dict[str, Any]], List[Dict[str, Any]]]:
    dirs: List[Dict[str, Any]] = []
    files: List[Dict[str, Any]] = []
    stats_cache: Dict[str, Dict[str, int]] = {}

    current_depth = path_depth_from_root(abs_dir, root_dir)

    for name in sorted(os.listdir(abs_dir), key=lambda x: x.lower()):
        full = os.path.join(abs_dir, name)
        try:
            st = os.stat(full)
        except Exception:
            continue

        item = {
            "name": name,
            "path": make_rel_from_root(full, root_dir),
            "mtime": datetime.fromtimestamp(st.st_mtime).strftime("%Y-%m-%d %H:%M:%S"),
        }

        if os.path.isdir(full):
            depth = path_depth_from_root(full, root_dir)
            if not matches_explorer_level_rule(depth, name):
                continue

            dir_info = build_dir_info(full, root_dir, stats_cache)
            item["type"] = "dir"
            item["depth"] = depth
            item["can_upload"] = depth == UPLOAD_ALLOWED_DEPTH
            item["has_children"] = False if depth >= EXPLORER_MAX_DEPTH else bool(dir_info.get("has_children"))
            item["expandable"] = depth < EXPLORER_MAX_DEPTH and bool(dir_info.get("has_children"))
            item["file_count"] = int(dir_info.get("file_count") or 0)
            item["total_size_bytes"] = int(dir_info.get("total_size_bytes") or 0)
            dirs.append(item)
        else:
            item["type"] = "file"
            item["size_bytes"] = st.st_size
            item["depth"] = current_depth
            files.append(item)

    return dirs, files


def sse_event(event: str, data: Dict[str, Any]) -> str:
    return f"event: {event}\ndata: {json.dumps(data, ensure_ascii=False)}\n\n"


def safe_err(msg: str) -> str:
    if not msg:
        return "不明なエラー"
    msg = re.sub(r"(app-[A-Za-z0-9_\-]{10,})", "app-***REDACTED***", msg)
    msg = re.sub(r"(Bearer\s+)[A-Za-z0-9_\-\.=]+", r"\1***REDACTED***", msg, flags=re.IGNORECASE)
    msg = re.sub(r"https?://[^\s]+", "[URL_REDACTED]", msg)
    return msg[:700]


# -----------------------
# Chat conversion (/chat-messages)
# -----------------------

def convert_via_dify_chat_messages_secure(
    api_base: str,
    api_key: str,
    user: str,
    source_path: str,
    source_meta: Dict[str, str],
    text: str,
    knowledge_style: str,
    chunk_sep: str,
) -> str:
    url = f"{api_base}/chat-messages"
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }

    instruction = build_rag_instruction(
        source_path=source_path,
        source_meta=source_meta,
        knowledge_style=knowledge_style,
        chunk_sep=chunk_sep,
    )

    query = (
        instruction
        + "\n\n===== SOURCE TEXT BEGIN =====\n"
        + text
        + "\n===== SOURCE TEXT END =====\n"
    )

    payload = {
        "inputs": {},
        "query": query,
        "response_mode": "blocking",
        "user": user,
    }

    try:
        r = requests.post(url, headers=headers, json=payload, timeout=REQ_TIMEOUT_SEC)
    except requests.RequestException:
        raise RuntimeError("API通信に失敗しました（ネットワーク/タイムアウト）。")

    if r.status_code >= 400:
        raise RuntimeError(f"APIエラー（HTTP {r.status_code}）: {safe_err(r.text)}")

    try:
        data = r.json()
    except Exception:
        raise RuntimeError("APIレスポンスの解析に失敗しました。")

    answer = data.get("answer")
    if not answer or not isinstance(answer, str):
        raise RuntimeError("APIレスポンスが想定外です（answerがありません）。")

    return answer.strip() + "\n"


def build_rag_instruction(source_path: str, source_meta: Dict[str, str], knowledge_style: str, chunk_sep: str) -> str:
    meta_lines = "\n".join([f"- {k}: {v}" for k, v in source_meta.items()])
    ext = (source_meta.get("ext") or "").lower()

    first_chunk_rule = f"""
        # 最初のチャンク（必須）
        - 出力の最初のチャンクは必ず「全体構成（目次/分類）」にする。
        - 形式例：
        - 見出し: 「## 全体構成（目次/分類）」
        - 次の1文: 「このチャンクでは文書全体の構成（目次）と分類方針を示す。」
        - 続けて、章立て（大カテゴリ）と、その中で扱う内容の要約を箇条書きで書く。
        - そのチャンクの末尾に必ず「{chunk_sep}」を単独行で置く。
        """

    excel_rules = ""
    if ext in {".xlsx", ".xls", ".xlsm"} and knowledge_style != "rag_natural":
        excel_rules = f"""
        # Excel特別ルール（標準/FAQ用）
        - 入力には [HEADER] と [ROW] が含まれる。
        - 出力は「データ行（[ROW]）1つにつき、必ずチャンク1つ」にする。
        - チャンク区切りは必ず「{chunk_sep}」の単独行にする。
        - [ROW]を統合しない。行同士をまとめない。
        """

    if knowledge_style == "rag_natural":
        style_block = f"""
        出力はMarkdownで「RAG向けMarkdown（自然言語）」として整形する。

        # 手順（必須）
        1) まず文書全体の構成を把握し、上位の章立て（大カテゴリ）を作る。
        2) 次に、人間が指示を出すような自然文で各チャンクの目的を宣言してから本文を置く。
        3) チャンク区切りは必ず「{chunk_sep}」の単独行にする。
        """
    elif knowledge_style == "faq":
        style_block = f"""
        出力はMarkdownで、FAQ形式にする。
        - 質問は具体的に、回答は短く「結論→根拠→例」の順にする。
        - チャンク区切りは必ず「{chunk_sep}」の単独行にする。
        """
    else:
        style_block = f"""
        出力はMarkdownで、RAGに最適化したナレッジへ整形する。
        - 文は「主語 + 述語」でできるだけ明確にする。
        - 検索されやすいキーワード（固有名詞/手順名/条件/例外/閾値）を含める。
        - チャンク区切りは必ず「{chunk_sep}」の単独行にする。
        - 情報を省略しない（重複は統合可）。
        """

    return f"""
        あなたは「社内RAG用ナレッジ整形AI」である。
        入力された文章を、検索精度が最大化するMarkdownへ再構成する。

        # 変換対象ファイル
        - path: {source_path}
        - meta:
        {meta_lines}

        # 絶対ルール
        - 出力は「変換後Markdown本文のみ」とする（前置き/解説/謝罪/注釈は禁止）。
        - 原文が曖昧な場合は「〜である可能性がある」等で補い、捏造しない。
        - チャンク区切りは必ず「{chunk_sep}」の単独行にする。

        {first_chunk_rule}

        {excel_rules}

        # スタイル
        {style_block}
        """.strip()


# -----------------------
# Markdown metadata + chunk analysis
# -----------------------

def _yaml_quote(v: str) -> str:
    s = "" if v is None else str(v)
    s = s.replace("\\", "\\\\").replace('"', '\\"')
    return f'"{s}"'


def attach_source_metadata(md: str, source_relpath: str, source_abspath: str, source_meta: Dict[str, str]) -> str:
    fm = {
        "source_relpath": source_relpath,
        "source_abspath": os.path.abspath(source_abspath),
        "source_filename": source_meta.get("filename") or os.path.basename(source_abspath),
        "source_ext": source_meta.get("ext") or os.path.splitext(source_abspath)[1].lower(),
        "source_size_bytes": source_meta.get("size_bytes") or "",
        "source_mtime": source_meta.get("mtime") or "",
        "generated_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
    }

    lines = ["---"]
    for k, v in fm.items():
        lines.append(f"{k}: {_yaml_quote(v)}")
    lines.append("---")
    lines.append("")

    body = (md or "").lstrip("\ufeff\n\r ")
    return "\n".join(lines) + body


def normalize_chunk_sep_lines(md: str, chunk_sep: str) -> str:
    sep = (chunk_sep or DEFAULT_CHUNK_SEP).strip() or DEFAULT_CHUNK_SEP
    lines = []
    for ln in (md or "").splitlines():
        if ln.strip() == sep:
            lines.append(sep)
        else:
            lines.append(ln.rstrip("\r"))
    out = "\n".join(lines).strip()
    return out + "\n"


def split_chunks(md: str, chunk_sep: str) -> List[str]:
    chunks: List[str] = []
    buf: List[str] = []
    sep = (chunk_sep or DEFAULT_CHUNK_SEP).strip()

    for ln in (md or "").splitlines():
        if ln.strip() == sep:
            txt = "\n".join(buf).strip()
            if txt:
                chunks.append(txt)
            buf = []
        else:
            buf.append(ln)

    last = "\n".join(buf).strip()
    if last:
        chunks.append(last)
    return chunks


def estimate_tokens(text: str) -> int:
    if not text:
        return 0
    total = len(text)
    if total <= 0:
        return 0

    ascii_cnt = sum(1 for ch in text if ord(ch) < 128)
    ascii_ratio = ascii_cnt / total

    chars_per_token = 3.0 if ascii_ratio >= 0.60 else 1.6
    est = int(total / chars_per_token) + 1
    return max(1, est)


def analyze_chunks_for_dify(markdown: str, chunk_sep: str) -> Dict[str, Any]:
    chunks = split_chunks(markdown, chunk_sep)
    lens = [estimate_tokens(c) for c in chunks] if chunks else []

    if not lens:
        return {
            "chunks": 0,
            "chunk_tokens_max": 0,
            "dify_max_tokens": min(1000, DIFY_MAX_SEG_TOKENS),
        }

    max_tok = max(lens)

    target = max_tok + 32
    target = max(200, target)
    target = min(DIFY_MAX_SEG_TOKENS, target)

    return {
        "chunks": len(chunks),
        "chunk_tokens_max": max_tok,
        "dify_max_tokens": target,
    }


# -----------------------
# Dify Knowledge API
# -----------------------

def dify_headers(api_key: str) -> Dict[str, str]:
    return {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }


def dify_list_datasets(api_base: str, api_key: str, prefix: str, limit: int = 100) -> List[Dict[str, str]]:
    out: List[Dict[str, str]] = []
    page = 1

    while True:
        url = f"{api_base}/datasets?page={page}&limit={limit}"
        r = requests.get(url, headers={"Authorization": f"Bearer {api_key}"}, timeout=REQ_TIMEOUT_SEC)
        if r.status_code >= 400:
            raise RuntimeError(f"datasets取得失敗（HTTP {r.status_code}）: {safe_err(r.text)}")

        data = r.json() if r.content else {}
        items = data.get("data") or []

        for it in items:
            name = (it.get("name") or "").strip()
            if prefix and not name.startswith(prefix):
                continue
            did = (it.get("id") or "").strip()
            if did and name:
                out.append({"id": did, "name": name})

        has_more = bool(data.get("has_more"))
        if not has_more:
            break

        page += 1
        if page > 200:
            break

    return out


def dify_get_dataset_detail(api_base: str, api_key: str, dataset_id: str) -> Dict[str, Any]:
    url = f"{api_base}/datasets/{dataset_id}"
    r = requests.get(url, headers={"Authorization": f"Bearer {api_key}"}, timeout=REQ_TIMEOUT_SEC)
    if r.status_code == 405:
        return {
            "id": dataset_id,
            "name": "",
            "_note": "datasets/{id} が GET 非対応のため、詳細は省略しました。",
        }
    if r.status_code >= 400:
        raise RuntimeError(f"dataset詳細取得失敗（HTTP {r.status_code}）: {safe_err(r.text)}")
    return r.json() if r.content else {}


def dify_list_documents_all(
    api_base: str,
    api_key: str,
    dataset_id: str,
    keyword: str = "",
    limit: int = 100,
) -> Tuple[List[Dict[str, Any]], int]:
    items_out: List[Dict[str, Any]] = []
    page = 1
    total = 0

    while True:
        qs = f"page={page}&limit={limit}"
        if keyword:
            qs += "&keyword=" + requests.utils.quote(keyword)

        url = f"{api_base}/datasets/{dataset_id}/documents?{qs}"
        r = requests.get(url, headers={"Authorization": f"Bearer {api_key}"}, timeout=REQ_TIMEOUT_SEC)
        if r.status_code >= 400:
            raise RuntimeError(f"documents取得失敗（HTTP {r.status_code}）: {safe_err(r.text)}")

        data = r.json() if r.content else {}
        items = data.get("data") or []
        total = int(data.get("total") or total or 0)

        for it in items:
            if isinstance(it, dict):
                items_out.append(it)

        if not bool(data.get("has_more")):
            break

        page += 1
        if page > 200:
            break

    return items_out, total


def dify_get_document_detail(
    api_base: str,
    api_key: str,
    dataset_id: str,
    document_id: str,
    metadata: str = "without",
) -> Dict[str, Any]:
    meta = metadata.strip() if metadata else "without"
    if meta not in {"all", "only", "without"}:
        meta = "without"

    url = f"{api_base}/datasets/{dataset_id}/documents/{document_id}?metadata={meta}"
    r = requests.get(url, headers={"Authorization": f"Bearer {api_key}"}, timeout=REQ_TIMEOUT_SEC)
    if r.status_code >= 400:
        raise RuntimeError(f"document詳細取得失敗（HTTP {r.status_code}）: {safe_err(r.text)}")
    return r.json() if r.content else {}


def dify_list_segments_page(
    api_base: str,
    api_key: str,
    dataset_id: str,
    document_id: str,
    page: int = 1,
    limit: int = 20,
    keyword: str = "",
    status: str = "",
) -> Dict[str, Any]:
    qs = f"page={page}&limit={limit}"
    if keyword:
        qs += "&keyword=" + requests.utils.quote(keyword)
    if status:
        qs += "&status=" + requests.utils.quote(status)

    url = f"{api_base}/datasets/{dataset_id}/documents/{document_id}/segments?{qs}"
    r = requests.get(url, headers={"Authorization": f"Bearer {api_key}"}, timeout=REQ_TIMEOUT_SEC)
    if r.status_code >= 400:
        raise RuntimeError(f"segments取得失敗（HTTP {r.status_code}）: {safe_err(r.text)}")

    data = r.json() if r.content else {}
    return {
        "items": data.get("data") or [],
        "has_more": bool(data.get("has_more")),
        "total": int(data.get("total") or 0),
        "page": int(data.get("page") or page),
        "limit": int(data.get("limit") or limit),
    }


def dify_get_segment_detail(
    api_base: str,
    api_key: str,
    dataset_id: str,
    document_id: str,
    segment_id: str,
) -> Dict[str, Any]:
    url = f"{api_base}/datasets/{dataset_id}/documents/{document_id}/segments/{segment_id}"
    r = requests.get(url, headers={"Authorization": f"Bearer {api_key}"}, timeout=REQ_TIMEOUT_SEC)
    if r.status_code >= 400:
        raise RuntimeError(f"segment詳細取得失敗（HTTP {r.status_code}）: {safe_err(r.text)}")
    return r.json() if r.content else {}


def dify_create_document_by_text(
    dataset_id: str,
    name: str,
    text: str,
    chunk_sep: str,
    dify_max_tokens: int,
    search_method: str = "hybrid_search",
) -> Tuple[str, str]:
    url = f"{DATASET_API_BASE}/datasets/{dataset_id}/document/create-by-text"

    payload: Dict[str, Any] = {
        "name": name,
        "text": text,
        "indexing_technique": "high_quality",
        "doc_form": "text_model",
        "process_rule": {
            "mode": "custom",
            "rules": {
                "pre_processing_rules": [
                    {"id": "remove_extra_spaces", "enabled": True},
                    {"id": "remove_urls_emails", "enabled": True},
                ],
                "segmentation": {
                    "separator": chunk_sep,
                    "max_tokens": int(dify_max_tokens),
                },
            },
        },
        "retrieval_model": {
            "search_method": search_method,
            "reranking_enable": False,
            "top_k": 5,
            "score_threshold_enabled": False,
        },
    }

    r = requests.post(url, headers=dify_headers(DATASET_API_KEY), json=payload, timeout=REQ_TIMEOUT_SEC)
    if r.status_code >= 400:
        raise RuntimeError(f"create-by-text 失敗（HTTP {r.status_code}）: {safe_err(r.text)}")

    data = r.json() if r.content else {}
    doc = data.get("document") or {}
    doc_id = (doc.get("id") or "").strip()
    batch = (data.get("batch") or "").strip()

    if not doc_id or not batch:
        raise RuntimeError("create-by-text レスポンスが想定外です（document.id / batch がありません）。")

    return doc_id, batch


def dify_get_indexing_status(dataset_id: str, batch: str) -> List[Dict[str, Any]]:
    url = f"{DATASET_API_BASE}/datasets/{dataset_id}/documents/{batch}/indexing-status"
    r = requests.get(url, headers={"Authorization": f"Bearer {DATASET_API_KEY}"}, timeout=REQ_TIMEOUT_SEC)
    if r.status_code >= 400:
        raise RuntimeError(f"indexing-status 取得失敗（HTTP {r.status_code}）: {safe_err(r.text)}")

    data = r.json() if r.content else {}
    return data.get("data") or []


def register_markdown_to_dify(dataset_id: str, doc_name: str, markdown: str, chunk_sep: str) -> Dict[str, Any]:
    stats = analyze_chunks_for_dify(markdown, chunk_sep)
    sep = (chunk_sep or DEFAULT_CHUNK_SEP).strip()

    doc_id, batch = dify_create_document_by_text(
        dataset_id=dataset_id,
        name=doc_name,
        text=markdown,
        chunk_sep=sep,
        dify_max_tokens=int(stats["dify_max_tokens"]),
        search_method="hybrid_search",
    )

    return {
        "doc_id": doc_id,
        "batch": batch,
        "chunk_sep": sep,
        "chunks": stats["chunks"],
        "chunk_tokens_max": stats["chunk_tokens_max"],
        "dify_max_tokens": stats["dify_max_tokens"],
        "search_method": "hybrid_search",
    }


def iter_indexing_status(dataset_id: str, batch: str, doc_id: str):
    start = time.time()
    last_key = None

    while True:
        if time.time() - start > INDEXING_MAX_WAIT_SEC:
            raise RuntimeError("ナレッジ埋め込みがタイムアウトしました。")

        items = dify_get_indexing_status(dataset_id, batch)

        target = None
        for it in items:
            if (it.get("id") or "").strip() == doc_id:
                target = it
                break

        if not target:
            time.sleep(INDEXING_POLL_SEC)
            continue

        st = (target.get("indexing_status") or "").strip()
        completed = int(target.get("completed_segments") or 0)
        total = int(target.get("total_segments") or 0)
        err = target.get("error")

        key = f"{st}:{completed}/{total}:{err}"
        if key != last_key:
            last_key = key
            terminal = st.lower() in {"completed", "error", "failed", "stopped"}
            yield {
                "indexing_status": st,
                "completed_segments": completed,
                "total_segments": total,
                "error": err,
                "terminal": terminal,
            }

            if terminal:
                return

        time.sleep(INDEXING_POLL_SEC)



_DATASET_CACHE_LOCK = threading.RLock()
_DATASET_CACHE: Dict[str, Any] = {
    "ts": 0.0,
    "items": [],
}

_DOCUMENT_CACHE_LOCK = threading.RLock()
_DOCUMENT_NAME_CACHE: Dict[str, Dict[str, Any]] = {}


def now_label() -> str:
    return datetime.now().strftime("%Y-%m-%d %H:%M:%S")


def normalize_name_key(name: str) -> str:
    return str(name or "").strip().lower()


def get_datasets_cached(force: bool = False) -> List[Dict[str, Any]]:
    with _DATASET_CACHE_LOCK:
        cache_ts = float(_DATASET_CACHE.get("ts") or 0.0)
        cache_items = list(_DATASET_CACHE.get("items") or [])
        fresh = (time.time() - cache_ts) <= ONDEMAND_DATASET_CACHE_TTL_SEC
        if cache_items and fresh and not force:
            return cache_items

    items = dify_list_datasets(
        api_base=DATASET_API_BASE,
        api_key=DATASET_API_KEY,
        prefix=DATASET_NAME_PREFIX,
        limit=100,
    )

    with _DATASET_CACHE_LOCK:
        _DATASET_CACHE["ts"] = time.time()
        _DATASET_CACHE["items"] = list(items or [])
        return list(_DATASET_CACHE["items"])


def find_dataset_by_name(dataset_name: str) -> Optional[Dict[str, Any]]:
    if not dataset_name:
        return None

    key = normalize_name_key(dataset_name)
    items = get_datasets_cached(force=False)
    for it in items:
        if normalize_name_key(it.get("name")) == key:
            return dict(it)
    return None


def get_dataset_document_name_keys_cached(dataset_id: str, force: bool = False) -> set:
    dataset_id = (dataset_id or "").strip()
    if not dataset_id:
        return set()

    with _DOCUMENT_CACHE_LOCK:
        entry = _DOCUMENT_NAME_CACHE.get(dataset_id) or {}
        cache_ts = float(entry.get("ts") or 0.0)
        cache_keys = set(entry.get("keys") or set())
        fresh = (time.time() - cache_ts) <= ONDEMAND_DOCUMENT_CACHE_TTL_SEC
        if cache_keys and fresh and not force:
            return cache_keys

    items, _ = dify_list_documents_all(
        api_base=DATASET_API_BASE,
        api_key=DATASET_API_KEY,
        dataset_id=dataset_id,
        keyword="",
        limit=100,
    )
    keys = {normalize_name_key((it or {}).get("name")) for it in (items or []) if (it or {}).get("name")}

    with _DOCUMENT_CACHE_LOCK:
        _DOCUMENT_NAME_CACHE[dataset_id] = {
            "ts": time.time(),
            "keys": set(keys),
        }
        return set(keys)


def dataset_document_exists_by_name(dataset_id: str, doc_name: str) -> bool:
    return normalize_name_key(doc_name) in get_dataset_document_name_keys_cached(dataset_id, force=False)


def remember_dataset_document_name(dataset_id: str, doc_name: str) -> None:
    dataset_id = (dataset_id or "").strip()
    if not dataset_id or not doc_name:
        return

    with _DOCUMENT_CACHE_LOCK:
        entry = _DOCUMENT_NAME_CACHE.get(dataset_id) or {"ts": time.time(), "keys": set()}
        keys = set(entry.get("keys") or set())
        keys.add(normalize_name_key(doc_name))
        entry["keys"] = keys
        entry["ts"] = time.time()
        _DOCUMENT_NAME_CACHE[dataset_id] = entry


def build_ondemand_dataset_name(rel_path: str) -> str:
    rel = (rel_path or "").strip().replace("\\", "/").strip("/")
    parts = [p for p in rel.split("/") if p]
    if len(parts) != UPLOAD_ALLOWED_DEPTH:
        return ""

    filtered = [p for p in parts if p != "元データ"]
    if not filtered:
        return ""

    return DATASET_NAME_PREFIX + "_".join(filtered)


def build_ondemand_markdown_path(folder_rel_path: str, original_name: str) -> Tuple[str, str, str]:
    rel = (folder_rel_path or "").strip().replace("\\", "/").strip("/")
    parts = [p for p in rel.split("/") if p]
    if len(parts) != UPLOAD_ALLOWED_DEPTH:
        raise RuntimeError("Lv5フォルダではないためMarkdown保存先を決定できません。")
    if len(parts) < 5 or parts[3] != "元データ":
        raise RuntimeError("元データ配下のLv5フォルダではないためMarkdown保存先を決定できません。")

    md_parts = list(parts)
    md_parts[3] = "マークダウン形式"
    md_dir_rel = "/".join(md_parts)
    md_dir_abs = resolve_explorer_path(EXPLORER_ROOT, md_dir_rel)

    base_name = os.path.splitext(sanitize_upload_filename(original_name))[0]
    if not base_name:
        raise RuntimeError("Markdownファイル名を決定できません。")

    md_name = f"{base_name}.md"
    md_abs_path = os.path.join(md_dir_abs, md_name)
    md_rel_path = make_rel_from_root(md_abs_path, EXPLORER_ROOT)
    return md_abs_path, md_rel_path, md_name


def is_ondemand_source_folder_rel(rel_path: str) -> bool:
    rel = (rel_path or "").strip().replace("\\", "/").strip("/")
    parts = [p for p in rel.split("/") if p]
    return len(parts) == UPLOAD_ALLOWED_DEPTH and len(parts) >= 5 and parts[3] == "元データ"


def strip_upload_timestamp_prefix(filename: str) -> str:
    safe = sanitize_upload_filename(filename)
    if not safe:
        return ""
    m = re.match(r"^\d{8}_\d{6}_(.+)$", safe)
    if m:
        return sanitize_upload_filename(m.group(1))
    return safe


def build_source_signature(source_abs_path: str, source_rel_path: str) -> str:
    rel = normalize_name_key((source_rel_path or "").replace("\\", "/"))
    try:
        st = os.stat(source_abs_path)
        size = int(st.st_size or 0)
        mtime_ns = int(getattr(st, "st_mtime_ns", int(float(st.st_mtime or 0) * 1_000_000_000)))
        return f"{rel}::{size}:{mtime_ns}"
    except Exception:
        return rel


def iter_ondemand_watch_folders(root_dir: str):
    root_abs = resolve_explorer_path(root_dir, "")

    for current_dir, dirnames, _ in os.walk(root_abs):
        depth = path_depth_from_root(current_dir, root_dir)

        if depth >= EXPLORER_MAX_DEPTH:
            dirnames[:] = []
        else:
            next_depth = depth + 1
            dirnames[:] = [d for d in dirnames if matches_explorer_level_rule(next_depth, d)]
            dirnames.sort(key=lambda x: x.lower())

        if depth == UPLOAD_ALLOWED_DEPTH:
            rel = make_rel_from_root(current_dir, root_dir)
            if is_ondemand_source_folder_rel(rel):
                yield current_dir, rel


def list_ondemand_source_files(folder_abs_path: str, root_dir: str) -> List[Dict[str, Any]]:
    out: List[Dict[str, Any]] = []
    for name in sorted(os.listdir(folder_abs_path), key=lambda x: x.lower()):
        abs_path = os.path.join(folder_abs_path, name)
        if not os.path.isfile(abs_path):
            continue

        ext = os.path.splitext(name)[1].lower()
        if ext not in ALLOWED_EXTS:
            continue

        rel_path = make_rel_from_root(abs_path, root_dir)
        original_name = strip_upload_timestamp_prefix(name) or sanitize_upload_filename(name)
        if not original_name:
            continue

        out.append({
            "source_abs_path": abs_path,
            "source_rel_path": rel_path,
            "source_saved_name": name,
            "source_original_name": original_name,
            "source_signature": build_source_signature(abs_path, rel_path),
        })
    return out


def build_ondemand_doc_key(dataset_id: str, dataset_name: str, markdown_name: str) -> str:
    ds_key = normalize_name_key(dataset_id or dataset_name)
    md_key = normalize_name_key(markdown_name)
    if not ds_key or not md_key:
        return ""
    return f"{ds_key}::{md_key}"


class OnDemandQueueManager:
    def __init__(self):
        self._lock = threading.RLock()
        self._cv = threading.Condition(self._lock)
        self._started = False
        self._thread: Optional[threading.Thread] = None
        self._tasks: Dict[str, Dict[str, Any]] = {}
        self._task_order: List[str] = []
        self._folder_queues: Dict[str, deque] = {}
        self._ready_folders: deque = deque()
        self._running_task_id: str = ""
        self._running_folder: str = ""
        self._active_doc_keys: Dict[str, str] = {}
        self._handled_source_signatures: Dict[str, Dict[str, Any]] = {}
        self._handled_source_order: deque = deque()

    def start(self) -> None:
        with self._lock:
            if self._started:
                return
            self._thread = threading.Thread(target=self._worker_loop, name="ondemand-queue-worker", daemon=True)
            self._thread.start()
            self._started = True

    def enqueue_saved_file(
        self,
        folder_rel_path: str,
        folder_abs_path: str,
        source_abs_path: str,
        source_rel_path: str,
        source_saved_name: str,
        source_original_name: str,
        source_signature: str = "",
        dataset_hint: Optional[Dict[str, Any]] = None,
        queue_message: str = "",
    ) -> Dict[str, Any]:
        folder_rel = (folder_rel_path or "").strip().replace("\\", "/").strip("/")
        source_signature = (source_signature or build_source_signature(source_abs_path, source_rel_path)).strip()
        dataset_name = build_ondemand_dataset_name(folder_rel)
        dataset = dict(dataset_hint) if isinstance(dataset_hint, dict) and dataset_hint else None
        message = ""

        if not API_BASE or not API_KEY:
            message = "生成AI API設定が未完了です。"
        elif not DATASET_API_BASE or not DATASET_API_KEY:
            message = "ナレッジAPI設定が未完了です。"
        elif not dataset_name:
            message = "ナレッジ名を判定できません。"
        elif not dataset:
            try:
                dataset = find_dataset_by_name(dataset_name)
            except Exception as e:
                dataset = None
                message = safe_err(str(e)) or "ナレッジ一覧の取得に失敗しました。"
            if not dataset and not message:
                message = "ナレッジが存在しません。管理者に問い合わせてください"

        md_abs_path = ""
        md_rel_path = ""
        md_name = ""
        try:
            md_abs_path, md_rel_path, md_name = build_ondemand_markdown_path(folder_rel, source_original_name)
        except Exception as e:
            if not message:
                message = safe_err(str(e))

        doc_key = build_ondemand_doc_key((dataset or {}).get("id") or "", dataset_name, md_name)

        with self._cv:
            handled = self._handled_source_signatures.get(source_signature) if source_signature else None
            if handled:
                return dict(handled.get("snapshot") or {})

            existing = self._find_task_by_source_signature_locked(source_signature)
            if existing:
                return self._public_task_snapshot(existing, self._queue_order_for_task_locked(existing.get("id") or ""))

            if doc_key:
                existing_task_id = self._active_doc_keys.get(doc_key) or ""
                existing_task = self._tasks.get(existing_task_id) if existing_task_id else None
                if existing_task:
                    return self._public_task_snapshot(existing_task, self._queue_order_for_task_locked(existing_task_id))

            task_id = uuid.uuid4().hex
            now = now_label()
            task = {
                "id": task_id,
                "folder_rel_path": folder_rel,
                "folder_abs_path": folder_abs_path,
                "folder_display": folder_rel or ".",
                "source_abs_path": source_abs_path,
                "source_rel_path": source_rel_path,
                "source_saved_name": source_saved_name,
                "source_original_name": source_original_name,
                "source_display_name": source_original_name or source_saved_name,
                "source_signature": source_signature,
                "dataset_name": dataset_name,
                "dataset_id": (dataset or {}).get("id") or "",
                "markdown_abs_path": md_abs_path,
                "markdown_rel_path": md_rel_path,
                "markdown_name": md_name,
                "doc_key": doc_key,
                "status": "queued" if dataset and md_abs_path else "error",
                "stage": "順番待ち" if dataset and md_abs_path else "受付不可",
                "message": (queue_message or "アップロード待ちキューに追加しました。") if dataset and md_abs_path else message,
                "attempt_no": 0,
                "retry_count": 0,
                "max_retry_count": ONDEMAND_QUEUE_MAX_RETRIES,
                "created_at": now,
                "updated_at": now,
                "started_at": "",
                "finished_at": "",
                "terminal": not bool(dataset and md_abs_path),
                "doc_id": "",
                "batch": "",
                "indexing_status": "",
                "completed_segments": 0,
                "total_segments": 0,
                "last_error": message if message else "",
                "markdown_written": False,
                "queue_order": None,
                "result": "pending" if dataset and md_abs_path else "error",
            }

            self._tasks[task_id] = task
            self._task_order.append(task_id)
            self._prune_locked()
            if task["status"] == "queued":
                if doc_key:
                    self._active_doc_keys[doc_key] = task_id
                fq = self._folder_queues.setdefault(folder_rel, deque())
                fq.append(task_id)
                self._ensure_folder_ready_locked(folder_rel)
                self._cv.notify_all()

            return self._public_task_snapshot(task, queue_order=None)

    def get_snapshot(self, limit: int = 200) -> Dict[str, Any]:
        with self._lock:
            queue_order_map = self._build_queue_order_map_locked()
            running_id = self._running_task_id
            items: List[Dict[str, Any]] = []
            summary = {
                "queued": 0,
                "running": 0,
                "completed": 0,
                "skipped": 0,
                "error": 0,
                "total": len(self._tasks),
            }

            for seq, task_id in enumerate(self._task_order, start=1):
                task = self._tasks.get(task_id)
                if not task:
                    continue
                st = str(task.get("status") or "")
                if st in summary:
                    summary[st] += 1
                qord = 0 if task_id == running_id else queue_order_map.get(task_id)
                item = self._public_task_snapshot(task, qord)
                item["_seq"] = seq
                items.append(item)

            items.sort(key=self._sort_key)
            for item in items:
                item.pop("_seq", None)
            if limit > 0:
                items = items[:limit]

            return {
                "summary": summary,
                "items": items,
            }

    def get_task_snapshot_by_id(self, task_id: str) -> Optional[Dict[str, Any]]:
        with self._lock:
            task = self._tasks.get(task_id)
            if not task:
                return None
            return self._public_task_snapshot(task, self._queue_order_for_task_locked(task_id))

    def get_task_snapshot_by_source_signature(self, source_signature: str) -> Optional[Dict[str, Any]]:
        sig = (source_signature or "").strip()
        if not sig:
            return None
        with self._lock:
            handled = self._handled_source_signatures.get(sig)
            if handled:
                return dict(handled.get("snapshot") or {})
            task = self._find_task_by_source_signature_locked(sig)
            if not task:
                return None
            return self._public_task_snapshot(task, self._queue_order_for_task_locked(task.get("id") or ""))

    def remember_handled_source_signature(self, source_signature: str, status: str, stage: str, message: str, result: str) -> None:
        snapshot = {
            "id": "",
            "folder_rel_path": "",
            "folder_display": "",
            "source_display_name": "",
            "source_saved_name": "",
            "source_rel_path": "",
            "dataset_name": "",
            "dataset_id": "",
            "markdown_name": "",
            "markdown_rel_path": "",
            "status": status,
            "stage": stage,
            "message": message,
            "attempt_no": 0,
            "retry_count": 0,
            "max_retry_count": ONDEMAND_QUEUE_MAX_RETRIES,
            "created_at": "",
            "updated_at": now_label(),
            "started_at": "",
            "finished_at": now_label(),
            "terminal": True,
            "doc_id": "",
            "batch": "",
            "indexing_status": "",
            "completed_segments": 0,
            "total_segments": 0,
            "queue_order": None,
            "last_error": "",
            "result": result,
        }
        with self._lock:
            self._remember_handled_source_signature_locked(source_signature, snapshot)

    def _sort_key(self, item: Dict[str, Any]):
        status = str(item.get("status") or "")
        queue_order = item.get("queue_order")
        seq = int(item.get("_seq") or 0)
        if status == "running":
            return (0, 0, 0)
        if status == "queued":
            return (1, int(queue_order or 999999), 0)
        return (2, 999999, -seq)

    def _public_task_snapshot(self, task: Dict[str, Any], queue_order: Optional[int]) -> Dict[str, Any]:
        return {
            "id": task.get("id") or "",
            "folder_rel_path": task.get("folder_rel_path") or "",
            "folder_display": task.get("folder_display") or "",
            "source_display_name": task.get("source_display_name") or "",
            "source_saved_name": task.get("source_saved_name") or "",
            "source_rel_path": task.get("source_rel_path") or "",
            "dataset_name": task.get("dataset_name") or "",
            "dataset_id": task.get("dataset_id") or "",
            "markdown_name": task.get("markdown_name") or "",
            "markdown_rel_path": task.get("markdown_rel_path") or "",
            "status": task.get("status") or "",
            "stage": task.get("stage") or "",
            "message": task.get("message") or "",
            "attempt_no": int(task.get("attempt_no") or 0),
            "retry_count": int(task.get("retry_count") or 0),
            "max_retry_count": int(task.get("max_retry_count") or 0),
            "created_at": task.get("created_at") or "",
            "updated_at": task.get("updated_at") or "",
            "started_at": task.get("started_at") or "",
            "finished_at": task.get("finished_at") or "",
            "terminal": bool(task.get("terminal")),
            "doc_id": task.get("doc_id") or "",
            "batch": task.get("batch") or "",
            "indexing_status": task.get("indexing_status") or "",
            "completed_segments": int(task.get("completed_segments") or 0),
            "total_segments": int(task.get("total_segments") or 0),
            "queue_order": queue_order if queue_order is not None else None,
            "last_error": task.get("last_error") or "",
            "result": task.get("result") or "",
        }

    def _prune_locked(self) -> None:
        if len(self._task_order) <= ONDEMAND_QUEUE_HISTORY_LIMIT:
            return

        removable = len(self._task_order) - ONDEMAND_QUEUE_HISTORY_LIMIT
        kept: List[str] = []
        for task_id in self._task_order:
            task = self._tasks.get(task_id)
            if not task:
                continue
            if removable > 0 and task.get("terminal"):
                self._tasks.pop(task_id, None)
                removable -= 1
                continue
            kept.append(task_id)
        self._task_order = kept

    def _ensure_folder_ready_locked(self, folder_rel_path: str) -> None:
        folder = (folder_rel_path or "").strip()
        if not folder:
            return
        if folder == self._running_folder:
            return
        if folder in self._ready_folders:
            return
        if self._folder_queues.get(folder):
            self._ready_folders.append(folder)

    def _build_queue_order_map_locked(self) -> Dict[str, int]:
        order_map: Dict[str, int] = {}
        temp_queues: Dict[str, deque] = {}
        for folder, q in self._folder_queues.items():
            temp_queues[folder] = deque(q)

        temp_ready = deque(self._ready_folders)
        running_task = self._tasks.get(self._running_task_id or "") if self._running_task_id else None
        if running_task:
            running_folder = str(running_task.get("folder_rel_path") or "")
            if temp_queues.get(running_folder):
                temp_ready.append(running_folder)

        order = 1
        while temp_ready:
            folder = temp_ready.popleft()
            q = temp_queues.get(folder)
            if not q:
                continue
            task_id = q.popleft()
            if task_id:
                order_map[task_id] = order
                order += 1
            if q:
                temp_ready.append(folder)
            else:
                temp_queues.pop(folder, None)

        return order_map

    def _queue_order_for_task_locked(self, task_id: str) -> Optional[int]:
        if not task_id:
            return None
        if task_id == self._running_task_id:
            return 0
        return self._build_queue_order_map_locked().get(task_id)

    def _find_task_by_source_signature_locked(self, source_signature: str) -> Optional[Dict[str, Any]]:
        sig = (source_signature or "").strip()
        if not sig:
            return None
        for task_id in self._task_order:
            task = self._tasks.get(task_id)
            if not task:
                continue
            if str(task.get("source_signature") or "") == sig:
                return task
        return None

    def _remember_handled_source_signature_locked(self, source_signature: str, snapshot: Dict[str, Any]) -> None:
        sig = (source_signature or "").strip()
        if not sig:
            return
        self._handled_source_signatures[sig] = {
            "snapshot": dict(snapshot or {}),
            "ts": time.time(),
        }
        self._handled_source_order.append(sig)
        while len(self._handled_source_order) > ONDEMAND_SEEN_SIGNATURE_LIMIT:
            old = self._handled_source_order.popleft()
            if old == sig:
                continue
            self._handled_source_signatures.pop(old, None)

    def _update_task(self, task_id: str, **fields: Any) -> Dict[str, Any]:
        with self._lock:
            task = self._tasks[task_id]
            for key, value in fields.items():
                task[key] = value
            task["updated_at"] = now_label()
            return dict(task)

    def _requeue_task_after_retry(self, task_id: str) -> None:
        with self._cv:
            task = self._tasks[task_id]
            folder = str(task.get("folder_rel_path") or "")
            fq = self._folder_queues.setdefault(folder, deque())
            fq.append(task_id)
            self._ensure_folder_ready_locked(folder)
            self._cv.notify_all()

    def _finish_task(self, task_id: str, status: str, stage: str, message: str, result: str = "") -> None:
        with self._lock:
            task = self._tasks[task_id]
            task["status"] = status
            task["stage"] = stage
            task["message"] = message
            task["terminal"] = True
            task["finished_at"] = now_label()
            task["updated_at"] = task["finished_at"]
            if result:
                task["result"] = result
            doc_key = str(task.get("doc_key") or "")
            if doc_key and self._active_doc_keys.get(doc_key) == task_id:
                self._active_doc_keys.pop(doc_key, None)
            if result in {"completed", "skipped"}:
                self._remember_handled_source_signature_locked(
                    str(task.get("source_signature") or ""),
                    self._public_task_snapshot(task, queue_order=None),
                )

    def _cleanup_markdown_if_needed(self, task_id: str) -> None:
        with self._lock:
            task = self._tasks[task_id]
            md_abs_path = str(task.get("markdown_abs_path") or "")
            md_written = bool(task.get("markdown_written"))
        if not md_abs_path or not md_written:
            return
        try:
            if os.path.exists(md_abs_path):
                os.remove(md_abs_path)
        except Exception:
            pass
        with self._lock:
            if task_id in self._tasks:
                self._tasks[task_id]["markdown_written"] = False

    def _worker_loop(self) -> None:
        while True:
            task_id = ""
            folder = ""
            with self._cv:
                while True:
                    while not self._ready_folders:
                        self._cv.wait(timeout=1.0)
                    folder = self._ready_folders.popleft()
                    q = self._folder_queues.get(folder)
                    if not q:
                        self._folder_queues.pop(folder, None)
                        continue
                    task_id = q.popleft()
                    if not q:
                        self._folder_queues.pop(folder, None)
                    break

                self._running_task_id = task_id
                self._running_folder = folder

            try:
                self._process_one_attempt(task_id)
            except Exception as e:
                self._finish_task(
                    task_id,
                    status="error",
                    stage="内部エラー",
                    message=safe_err(str(e)),
                    result="error",
                )
                self._cleanup_markdown_if_needed(task_id)
            finally:
                with self._cv:
                    self._running_task_id = ""
                    self._running_folder = ""
                    if self._folder_queues.get(folder):
                        self._ensure_folder_ready_locked(folder)
                        self._cv.notify_all()

    def _process_one_attempt(self, task_id: str) -> None:
        with self._lock:
            task = self._tasks[task_id]
            task["attempt_no"] = int(task.get("attempt_no") or 0) + 1
            attempt_no = int(task["attempt_no"])
            task["status"] = "running"
            task["stage"] = "処理中"
            task["message"] = f"処理を開始しました（{attempt_no}回目）。"
            if not task.get("started_at"):
                task["started_at"] = now_label()
            task["updated_at"] = now_label()

        try:
            dataset_id = str(task.get("dataset_id") or "")
            source_abs_path = str(task.get("source_abs_path") or "")
            source_display_name = str(task.get("source_display_name") or "")
            markdown_name = str(task.get("markdown_name") or "")
            markdown_abs_path = str(task.get("markdown_abs_path") or "")
            folder_rel_path = str(task.get("folder_rel_path") or "")

            self._update_task(task_id, stage="差分確認", message="同一Markdownがナレッジに存在するか確認しています。")
            if dataset_document_exists_by_name(dataset_id, markdown_name):
                self._finish_task(
                    task_id,
                    status="skipped",
                    stage="差分なし",
                    message="同一Markdownが既にナレッジに存在するため登録をスキップしました。",
                    result="skipped",
                )
                return

            self._update_task(task_id, stage="テキスト抽出", message=f"{source_display_name} からテキストを抽出しています。")
            raw_text, meta = extract_text(source_abs_path, knowledge_style=ONDEMAND_QUEUE_STYLE)
            if not raw_text.strip():
                raise RuntimeError("抽出テキストが空でした。")
            if len(raw_text) > MAX_INPUT_CHARS:
                raw_text = raw_text[:MAX_INPUT_CHARS] + "\n...(truncated)\n"

            self._update_task(task_id, stage="Markdown変換", message="RAG向けMarkdownへ変換しています。")
            md_body = convert_via_dify_chat_messages_secure(
                api_base=API_BASE,
                api_key=API_KEY,
                user=ONDEMAND_QUEUE_USER,
                source_path=folder_rel_path + "/" + source_display_name if folder_rel_path else source_display_name,
                source_meta=meta,
                text=raw_text,
                knowledge_style=ONDEMAND_QUEUE_STYLE,
                chunk_sep=ONDEMAND_QUEUE_CHUNK_SEP,
            )
            md_body = normalize_chunk_sep_lines(md_body, ONDEMAND_QUEUE_CHUNK_SEP)
            md_save = attach_source_metadata(
                md_body,
                source_relpath=str(task.get("source_rel_path") or source_display_name),
                source_abspath=source_abs_path,
                source_meta=meta,
            )

            self._update_task(task_id, stage="Markdown保存", message="Markdownファイルを保存しています。")
            os.makedirs(os.path.dirname(markdown_abs_path), exist_ok=True)
            with open(markdown_abs_path, "w", encoding="utf-8", newline="\n") as f:
                f.write(md_save)
            self._update_task(task_id, markdown_written=True)

            self._update_task(task_id, stage="ナレッジ登録", message="Difyナレッジへ登録しています。")
            reg = register_markdown_to_dify(
                dataset_id=dataset_id,
                doc_name=markdown_name,
                markdown=md_body,
                chunk_sep=ONDEMAND_QUEUE_CHUNK_SEP,
            )
            self._update_task(
                task_id,
                doc_id=reg.get("doc_id") or "",
                batch=reg.get("batch") or "",
                message=(
                    f"受付済み: chunks={reg.get('chunks')} / max_tokens={reg.get('dify_max_tokens')} / search=hybrid_search"
                ),
            )

            final = None
            for prog in iter_indexing_status(dataset_id=dataset_id, batch=reg["batch"], doc_id=reg["doc_id"]):
                completed = int(prog.get("completed_segments") or 0)
                total = int(prog.get("total_segments") or 0)
                status = str(prog.get("indexing_status") or "")
                msg = f"埋め込み中: status={status} / segments={completed}/{total}"
                if prog.get("error"):
                    msg += f" / error={safe_err(str(prog.get('error')))}"
                self._update_task(
                    task_id,
                    stage="埋め込み待ち",
                    indexing_status=status,
                    completed_segments=completed,
                    total_segments=total,
                    message=msg,
                )
                if prog.get("terminal"):
                    final = prog
                    break

            if not final:
                raise RuntimeError("Dify埋め込みの進捗取得に失敗しました。")
            if str(final.get("indexing_status") or "").lower() != "completed":
                raise RuntimeError(
                    f"Dify埋め込み失敗: status={final.get('indexing_status')} error={final.get('error')}"
                )
            if int(final.get("total_segments") or 0) <= 0:
                raise RuntimeError("Dify側で0セグメントのまま完了しました。")

            remember_dataset_document_name(dataset_id, markdown_name)
            self._finish_task(
                task_id,
                status="completed",
                stage="完了",
                message="Markdown保存とナレッジ登録が完了しました。",
                result="completed",
            )
        except Exception as e:
            err = safe_err(str(e))
            with self._lock:
                task = self._tasks[task_id]
                retry_count = int(task.get("retry_count") or 0)
                max_retry = int(task.get("max_retry_count") or 0)
                task["last_error"] = err

            if retry_count < max_retry:
                retry_count += 1
                with self._lock:
                    task = self._tasks[task_id]
                    task["retry_count"] = retry_count
                    task["status"] = "queued"
                    task["stage"] = "リトライ待ち"
                    task["message"] = f"{err} / リトライ {retry_count}/{max_retry} を待機しています。"
                    task["updated_at"] = now_label()
                    task["terminal"] = False
                    task["result"] = "pending"
                self._requeue_task_after_retry(task_id)
                return

            self._cleanup_markdown_if_needed(task_id)
            self._finish_task(
                task_id,
                status="error",
                stage="エラー終了",
                message=f"{err} / リトライ上限に達したため中止しました。",
                result="error",
            )


class OnDemandFolderMonitor:
    def __init__(self, queue_manager: OnDemandQueueManager):
        self._queue = queue_manager
        self._lock = threading.RLock()
        self._started = False
        self._thread: Optional[threading.Thread] = None
        self._running = False
        self._last_scan_started_at = ""
        self._last_scan_finished_at = ""
        self._last_scan_error = ""
        self._last_stats: Dict[str, Any] = {
            "folders": 0,
            "files": 0,
            "enqueued": 0,
            "known": 0,
            "doc_exists": 0,
            "dataset_missing": 0,
            "not_target": 0,
        }

    def start(self) -> None:
        if not ONDEMAND_MONITOR_ENABLED:
            return
        with self._lock:
            if self._started:
                return
            self._thread = threading.Thread(target=self._loop, name="ondemand-folder-monitor", daemon=True)
            self._thread.start()
            self._started = True

    def get_status(self) -> Dict[str, Any]:
        with self._lock:
            return {
                "enabled": bool(ONDEMAND_MONITOR_ENABLED),
                "running": bool(self._running),
                "interval_sec": ONDEMAND_MONITOR_INTERVAL_SEC,
                "last_scan_started_at": self._last_scan_started_at,
                "last_scan_finished_at": self._last_scan_finished_at,
                "last_scan_error": self._last_scan_error,
                "last_stats": dict(self._last_stats),
            }

    def _set_scan_state(self, **fields: Any) -> None:
        with self._lock:
            for key, value in fields.items():
                setattr(self, key, value)

    def _loop(self) -> None:
        while True:
            started_at = now_label()
            self._set_scan_state(_running=True, _last_scan_started_at=started_at, _last_scan_error="")
            stats = {
                "folders": 0,
                "files": 0,
                "enqueued": 0,
                "known": 0,
                "doc_exists": 0,
                "dataset_missing": 0,
                "not_target": 0,
            }
            last_error = ""

            try:
                self._scan_once(stats)
            except Exception as e:
                last_error = safe_err(str(e))

            finished_at = now_label()
            with self._lock:
                self._running = False
                self._last_scan_finished_at = finished_at
                self._last_scan_error = last_error
                self._last_stats = dict(stats)

            time.sleep(ONDEMAND_MONITOR_INTERVAL_SEC)

    def _scan_once(self, stats: Dict[str, int]) -> None:
        if not API_BASE or not API_KEY or not DATASET_API_BASE or not DATASET_API_KEY:
            return

        datasets = get_datasets_cached(force=False)
        dataset_map = {normalize_name_key((it or {}).get("name")): dict(it) for it in (datasets or []) if (it or {}).get("name")}

        for folder_abs_path, folder_rel_path in iter_ondemand_watch_folders(EXPLORER_ROOT):
            stats["folders"] += 1

            if not is_ondemand_source_folder_rel(folder_rel_path):
                stats["not_target"] += 1
                continue

            dataset_name = build_ondemand_dataset_name(folder_rel_path)
            dataset = dataset_map.get(normalize_name_key(dataset_name)) if dataset_name else None

            files = list_ondemand_source_files(folder_abs_path, EXPLORER_ROOT)
            for info in files:
                stats["files"] += 1
                source_signature = str(info.get("source_signature") or "")
                if source_signature and self._queue.get_task_snapshot_by_source_signature(source_signature):
                    stats["known"] += 1
                    continue

                if not dataset:
                    stats["dataset_missing"] += 1
                    continue

                md_name = ""
                try:
                    _, _, md_name = build_ondemand_markdown_path(folder_rel_path, str(info.get("source_original_name") or ""))
                except Exception:
                    continue

                if dataset_document_exists_by_name(str(dataset.get("id") or ""), md_name):
                    self._queue.remember_handled_source_signature(
                        source_signature=source_signature,
                        status="skipped",
                        stage="差分なし",
                        message="同一Markdownが既にナレッジに存在するため監視対象から除外しました。",
                        result="skipped",
                    )
                    stats["doc_exists"] += 1
                    continue

                task = self._queue.enqueue_saved_file(
                    folder_rel_path=folder_rel_path,
                    folder_abs_path=folder_abs_path,
                    source_abs_path=str(info.get("source_abs_path") or ""),
                    source_rel_path=str(info.get("source_rel_path") or ""),
                    source_saved_name=str(info.get("source_saved_name") or ""),
                    source_original_name=str(info.get("source_original_name") or ""),
                    source_signature=source_signature,
                    dataset_hint=dataset,
                    queue_message="フォルダ監視で検出し、順番待ちキューへ追加しました。",
                )
                if task and str(task.get("status") or "") == "queued":
                    stats["enqueued"] += 1
                else:
                    stats["known"] += 1


ONDEMAND_QUEUE = OnDemandQueueManager()
ONDEMAND_MONITOR = OnDemandFolderMonitor(ONDEMAND_QUEUE)

if __name__ == "__main__":
    app = create_app()
    app.run(host="0.0.0.0", port=5212, debug=False, threaded=True)